import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# Define color palette
DARK_BLUE = RGBColor(11, 29, 58)
LIGHT_BG = RGBColor(245, 247, 250)
TEXT_DARK = RGBColor(33, 37, 41)
ACCENT_BLUE = RGBColor(30, 144, 255)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

def add_title(slide, text, color=DARK_BLUE):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Arial"

# ==========================================
# SLIDE 1: Title Slide (Dark Background Layout)
# ==========================================
slide_layout = prs.slide_layouts[6] # blank layout
slide1 = prs.slides.add_slide(slide_layout)
slide1.background.fill.solid()
slide1.background.fill.fore_color.rgb = DARK_BLUE

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Design and Implementation of an Intelligent RAG Chatbot for Ghana Revenue Authority Compliance"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.33), Inches(2.0))
tf_sub = sub_box.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Technical Stack: LangChain | FAISS | Hugging Face | Groq (Llama 3.1) | Streamlit\nPresenter: Lukeman Abubakar\nTrack: Advanced ML & AI"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = RGBColor(200, 210, 230)

# Data for remaining slides
slides_data = [
    {
        "title": "The Problem Statement",
        "bullets": [
            "Dense & Long: Official texts like GRA tax guides, VAT rules, and small business obligations span dozens of dense pages.",
            "Impractical to Read: For the average taxpayer or business owner, finding specific rules inside an unstructured PDF is highly impractical.",
            "The LLM Dilemma: Standard AI models lack access to local Ghanaian regulations and confidently fabricate or hallucinate answers, creating legal risks."
        ]
    },
    {
        "title": "The Solution: Retrieval-Augmented Generation (RAG)",
        "bullets": [
            "Eliminating Hallucinations: Instead of relying on an LLM's general training, RAG limits the model's pool of facts strictly to a verified document corpus.",
            "How It Works: The system searches the GRA documents first, pulls out relevant text snippets, and forces the AI to answer using ONLY that context.",
            "Trust & Verification: The system provides explicit source citations alongside answers, allowing users to verify facts themselves."
        ]
    },
    {
        "title": "System Architecture (Stage 1): Ingestion & Indexing",
        "bullets": [
            "Document Loading: Programmatically extracted raw text from official GRA compliance PDFs using LangChain loaders.",
            "Semantic Chunking: Split long files into smaller text blocks of 500 to 800 characters with text overlap to ensure no tax clause was cut in half mid-sentence.",
            "Vector Embeddings: Converted text chunks into dense vectors using Hugging Face's all-MiniLM-L6-v2 model.",
            "Local Index Storage: Indexed and stored vectors locally on disk using a highly efficient FAISS vector database."
        ]
    },
    {
        "title": "System Architecture (Stage 2): Retrieval & Chat Interface",
        "bullets": [
            "User Query Processing: The user enters a question in the Streamlit web interface; the question is immediately transformed into a vector.",
            "Semantic Search: FAISS executes a cosine-similarity calculation to instantly pull the top k=3 most relevant tax chunks.",
            "Groq Cloud Inference: Sends the text chunks and the user's question to a high-speed llama-3.1-8b-instant model on Groq.",
            "Response Generation: The UI presents a natural language answer grounded completely in the document context."
        ]
    },
    {
        "title": "Prompt Engineering & Guardrails",
        "bullets": [
            "Strict Context Constraint: System prompt commands the model to use ONLY the provided text.",
            "Out-of-Scope Rule: The AI is strictly ordered to decline general knowledge or unrelated questions.",
            "Mandated Fallback Text: If information is missing from the corpus, the model must output exactly: 'I could not find an answer to this question in the documents. You may want to consult the source document directly...'"
        ]
    },
    {
        "title": "Software Engineering Best Practices",
        "bullets": [
            "Credential Isolation: The Groq API authentication key is completely removed from the codebase and managed via a secure local .env file.",
            "Leaking Prevention: Formulated a rigorous .gitignore file to ensure virtual environments (.venv/) and secret files are completely hidden from public repositories.",
            "Reproducibility: Used 'pip freeze > requirements.txt' to capture the entire system dependency tree for easy deployment and evaluation by grading instructors."
        ]
    },
    {
        "title": "System Evaluation & Results",
        "bullets": [
            "Direct Fact Accuracy: Flawlessly handled specific lookups (e.g., late filing penalties, VAT registration thresholds) by isolating the exact matching text chunks.",
            "Information Synthesis: Successfully combined separate text chunks to answer complex queries (e.g., tracking simultaneous late filing fees and accumulation of interest).",
            "Guardrail Performance: 100% success rate intercepting irrelevant inputs (e.g., general geography, recipes) and triggering the mandatory fallback message."
        ]
    },
    {
        "title": "Operational Failure Modes",
        "bullets": [
            "Context Splitting: Fixed-character chunking can split a complex clause across borders, separating a core tax penalty from its legal exemptions.",
            "Free Tier API Rate Limits: Rapid testing of complex queries occasionally triggers HTTP 429 exceptions on free Hugging Face or Groq endpoints."
        ]
    },
    {
        "title": "Permanent System Limitations",
        "bullets": [
            "Document Currency: The chatbot is structurally dependent on static files; updates to GRA tax policies require a manual database rebuild.",
            "Language Barriers: The pipeline operates entirely in English, creating an accessibility bottleneck for local traders who speak Twi or Hausa.",
            "Professional Advice Disclaimer: The system acts as a document text search tool; it does not replace a certified tax professional."
        ]
    },
    {
        "title": "Summary & Future Outlook",
        "bullets": [
            "Successful Objectives: Built a functional, secure, end-to-end RAG application that provides transparent source citations.",
            "Future Work: Upgrade to hierarchical parent-child chunking to preserve split paragraphs, and integrate translation APIs to support queries in local Ghanaian languages.",
            "Thank You! Questions or feedback? (Enquiries: programsupport@thriveafrica.co)"
        ]
    }
]

# Generate standard content slides
for sdata in slides_data:
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, sdata["title"])
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for idx, bullet in enumerate(sdata["bullets"]):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK
        p.font.name = "Arial"
        p.space_after = Pt(24)

# Save presentation
prs.save("GRA_RAG_Capstone_Presentation.pptx")
print("PowerPoint presentation successfully created as 'GRA_RAG_Capstone_Presentation.pptx'!")